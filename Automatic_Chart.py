import pandas as pd
import plotly.express as px
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

def automatic_chart_analysis(df: pd.DataFrame):
    """Automatically generate charts based on the analysis of the DataFrame."""
    charts = []
    numerical_cols = df.select_dtypes(include=['number']).columns.tolist()
    categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()

    if numerical_cols and categorical_cols:
        for cat_col in categorical_cols:
            for num_col in numerical_cols:
                fig = px.bar(df, x=cat_col, y=num_col, title=f"{cat_col} vs {num_col}")
                charts.append(fig)
    for num_col in numerical_cols:
        fig = px.histogram(df, x=num_col, title=f"Histogram of {num_col}")
        charts.append(fig)
    for i in range(len(numerical_cols)):
        for j in range(i + 1, len(numerical_cols)):
            fig = px.scatter(df, x=numerical_cols[i], y=numerical_cols[j], title=f"{numerical_cols[i]} vs {numerical_cols[j]}")
            charts.append(fig)
    if len(numerical_cols) > 1:
        correlation_matrix = df[numerical_cols].corr()
        fig = px.imshow(correlation_matrix, text_auto=True, title="Correlation Matrix")
        charts.append(fig)
    return charts

def create_ppt(charts: list, df: pd.DataFrame, ppt_name: str = "Chart_Analysis") -> str:
    """Creates a PowerPoint presentation with chart summaries and data tables."""
    prs = Presentation()
    for idx, fig in enumerate(charts):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"Chart {idx + 1}: {fig.layout.title.text or 'No Title'}"

        # Add description box
        left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = (
            "Interactive chart available in the app or as HTML download.\n"
            "Use the downloaded CSV to recreate this chart in PowerPoint."
        )

        # Add data table (simplified example)
        if fig.data[0].x is not None and fig.data[0].y is not None:
            rows = min(len(fig.data[0].x), 5)  # Limit to 5 rows for brevity
            cols = 2
            table = slide.shapes.add_table(rows + 1, cols, Inches(1), Inches(4), Inches(8), Inches(2)).table
            table.cell(0, 0).text = fig.layout.xaxis.title.text or "X"
            table.cell(0, 1).text = fig.layout.yaxis.title.text or "Y"
            for i in range(rows):
                table.cell(i + 1, 0).text = str(fig.data[0].x[i])
                table.cell(i + 1, 1).text = str(fig.data[0].y[i])

    # Save to temporary file with dynamic name
    pptx_file = f"PPT_by_Hindav_{ppt_name}.pptx"
    prs.save(pptx_file)
    return pptx_file